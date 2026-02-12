import os
import tempfile
import re
from typing import List, Optional
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
from bs4 import BeautifulSoup
import requests
from youtube_transcript_api import YouTubeTranscriptApi


def extract_youtube_video_id(url: str) -> Optional[str]:
    patterns = [
        r'(?:v=|\/)([0-9A-Za-z_-]{11}).*',
        r'(?:embed\/)([0-9A-Za-z_-]{11})',
        r'(?:youtu\.be\/)([0-9A-Za-z_-]{11})',
    ]
    for pattern in patterns:
        match = re.search(pattern, url)
        if match:
            return match.group(1)
    return None


def load_pdf(file_path: str) -> List[Document]:
    reader = PdfReader(file_path)
    documents = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            documents.append(Document(
                page_content=text,
                metadata={"source": file_path, "page": i + 1, "type": "file"}
            ))
    if not documents:
        raise ValueError("PDF contains no extractable text")
    return documents


def load_docx(file_path: str) -> List[Document]:
    doc = DocxDocument(file_path)
    text = "\n".join([para.text for para in doc.paragraphs if para.text])
    if text and text.strip():
        return [Document(
            page_content=text,
            metadata={"source": file_path, "type": "file"}
        )]
    raise ValueError("Document contains no extractable text")


def load_pptx(file_path: str) -> List[Document]:
    prs = Presentation(file_path)
    documents = []
    for i, slide in enumerate(prs.slides):
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text)
        if text_parts:
            documents.append(Document(
                page_content="\n".join(text_parts),
                metadata={"source": file_path, "slide": i + 1, "type": "file"}
            ))
    return documents


def load_xlsx(file_path: str) -> List[Document]:
    wb = load_workbook(file_path, data_only=True)
    documents = []
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        rows = []
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join([str(cell) for cell in row if cell is not None])
            if row_text.strip():
                rows.append(row_text)
        if rows:
            documents.append(Document(
                page_content="\n".join(rows),
                metadata={"source": file_path, "sheet": sheet_name, "type": "file"}
            ))
    return documents


def get_base_domain(url: str) -> str:
    """Extract base domain from URL for same-domain link filtering."""
    from urllib.parse import urlparse
    parsed = urlparse(url)
    return f"{parsed.scheme}://{parsed.netloc}"


def normalize_url(base_url: str, href: str) -> Optional[str]:
    """Convert relative URLs to absolute and normalize."""
    from urllib.parse import urljoin, urlparse
    
    if not href or href.startswith(('#', 'javascript:', 'mailto:', 'tel:')):
        return None
    
    absolute_url = urljoin(base_url, href)
    parsed = urlparse(absolute_url)
    
    normalized = f"{parsed.scheme}://{parsed.netloc}{parsed.path}"
    if normalized.endswith('/'):
        normalized = normalized[:-1]
    
    return normalized


def extract_page_content(url: str, headers: dict) -> tuple:
    """Extract text content and links from a single page (simple HTTP request)."""
    try:
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        
        content_type = response.headers.get('Content-Type', '')
        if 'text/html' not in content_type:
            return None, []
        
        soup = BeautifulSoup(response.text, "html.parser")
        
        links = []
        for a_tag in soup.find_all('a', href=True):
            normalized = normalize_url(url, a_tag['href'])
            if normalized:
                links.append(normalized)
        
        for script in soup(["script", "style", "nav", "footer", "header", "aside"]):
            script.decompose()
        
        title = soup.title.string if soup.title else url
        text = soup.get_text(separator="\n", strip=True)
        text = re.sub(r'\n\s*\n', '\n\n', text)
        
        return (text, title), links
    except Exception:
        return None, []


def get_selenium_driver():
    """Create and return a Selenium WebDriver with Chrome (memory optimized)."""
    from selenium import webdriver
    from selenium.webdriver.chrome.options import Options
    import shutil
    
    chrome_options = Options()
    chrome_options.add_argument("--headless=new")
    chrome_options.add_argument("--no-sandbox")
    chrome_options.add_argument("--disable-dev-shm-usage")
    chrome_options.add_argument("--disable-gpu")
    chrome_options.add_argument("--disable-extensions")
    chrome_options.add_argument("--disable-plugins")
    chrome_options.add_argument("--disable-images")
    chrome_options.add_argument("--blink-settings=imagesEnabled=false")
    chrome_options.add_argument("--window-size=1280,720")
    chrome_options.add_argument("--single-process")
    chrome_options.add_argument("--disable-background-networking")
    chrome_options.add_argument("--disable-default-apps")
    chrome_options.add_argument("--disable-sync")
    chrome_options.add_argument("--js-flags=--max-old-space-size=256")
    chrome_options.add_argument("user-agent=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36")
    
    chromium_path = shutil.which("chromium")
    if chromium_path:
        chrome_options.binary_location = chromium_path
    
    driver = webdriver.Chrome(options=chrome_options)
    driver.set_page_load_timeout(15)
    driver.set_script_timeout(10)
    return driver


def extract_page_content_selenium(driver, url: str) -> tuple:
    """Extract text content and links from a page using Selenium (JavaScript rendered)."""
    import time
    try:
        driver.get(url)
        time.sleep(2)
        
        html = driver.page_source
        soup = BeautifulSoup(html, "html.parser")
        
        links = []
        for a_tag in soup.find_all('a', href=True):
            normalized = normalize_url(url, a_tag['href'])
            if normalized:
                links.append(normalized)
        
        for script in soup(["script", "style", "nav", "footer", "header", "aside", "noscript"]):
            script.decompose()
        
        title = soup.title.string if soup.title else url
        text = soup.get_text(separator="\n", strip=True)
        text = re.sub(r'\n\s*\n', '\n\n', text)
        text = re.sub(r'로딩 중\.\.\.', '', text)
        
        return (text, title), links
    except Exception as e:
        print(f"Error extracting {url}: {e}")
        return None, []


def crawl_website(start_url: str, max_pages: int = 30, progress_callback=None, use_js: bool = True) -> List[Document]:
    """Crawl entire website starting from the given URL.
    
    Args:
        start_url: The starting URL to crawl
        max_pages: Maximum number of pages to crawl
        progress_callback: Optional callback function(current, total, url)
        use_js: If True, use Selenium for JavaScript rendering (recommended for dynamic sites)
    """
    base_domain = get_base_domain(start_url)
    
    visited = set()
    to_visit = [start_url]
    documents = []
    consecutive_errors = 0
    max_consecutive_errors = 5
    driver_restart_interval = 10
    
    if use_js:
        driver = None
        try:
            driver = get_selenium_driver()
            pages_since_restart = 0
            
            while to_visit and len(visited) < max_pages:
                current_url = to_visit.pop(0)
                
                normalized_current = normalize_url(current_url, current_url)
                if normalized_current in visited:
                    continue
                
                visited.add(normalized_current)
                
                if progress_callback:
                    progress_callback(len(visited), max_pages, current_url)
                
                try:
                    if pages_since_restart >= driver_restart_interval:
                        try:
                            driver.quit()
                        except:
                            pass
                        driver = get_selenium_driver()
                        pages_since_restart = 0
                    
                    result, links = extract_page_content_selenium(driver, current_url)
                    pages_since_restart += 1
                    
                    if result:
                        text, title = result
                        if text and len(text.strip()) > 100:
                            documents.append(Document(
                                page_content=text,
                                metadata={
                                    "source": current_url,
                                    "title": title,
                                    "type": "website",
                                    "filename": title[:50] if title else current_url
                                }
                            ))
                        consecutive_errors = 0
                    else:
                        consecutive_errors += 1
                    
                    for link in links:
                        if link.startswith(base_domain) and link not in visited:
                            if not any(ext in link.lower() for ext in ['.pdf', '.jpg', '.png', '.gif', '.zip', '.mp4', '.mp3', '.hwp']):
                                if link not in to_visit:
                                    to_visit.append(link)
                                    
                except Exception as page_error:
                    print(f"Error on page {current_url}: {page_error}")
                    consecutive_errors += 1
                    
                    if consecutive_errors >= max_consecutive_errors:
                        print(f"Too many consecutive errors ({consecutive_errors}), restarting driver...")
                        try:
                            driver.quit()
                        except:
                            pass
                        driver = get_selenium_driver()
                        consecutive_errors = 0
                        pages_since_restart = 0
            
            if driver:
                try:
                    driver.quit()
                except:
                    pass
                    
        except Exception as e:
            print(f"Selenium error ({e}), falling back to simple HTTP requests")
            if driver:
                try:
                    driver.quit()
                except:
                    pass
            return crawl_website(start_url, max_pages, progress_callback, use_js=False)
    else:
        headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
        
        while to_visit and len(visited) < max_pages:
            current_url = to_visit.pop(0)
            
            normalized_current = normalize_url(current_url, current_url)
            if normalized_current in visited:
                continue
            
            visited.add(normalized_current)
            
            if progress_callback:
                progress_callback(len(visited), max_pages, current_url)
            
            result, links = extract_page_content(current_url, headers)
            
            if result:
                text, title = result
                if text and len(text.strip()) > 100:
                    documents.append(Document(
                        page_content=text,
                        metadata={
                            "source": current_url,
                            "title": title,
                            "type": "website",
                            "filename": title[:50] if title else current_url
                        }
                    ))
            
            for link in links:
                if link.startswith(base_domain) and link not in visited:
                    if not any(ext in link.lower() for ext in ['.pdf', '.jpg', '.png', '.gif', '.zip', '.mp4', '.mp3', '.hwp']):
                        if link not in to_visit:
                            to_visit.append(link)
    
    return documents


def load_website(url: str) -> List[Document]:
    """Load a single webpage (legacy function for backward compatibility)."""
    try:
        headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, "html.parser")
        
        for script in soup(["script", "style", "nav", "footer", "header"]):
            script.decompose()
        
        text = soup.get_text(separator="\n", strip=True)
        text = re.sub(r'\n\s*\n', '\n\n', text)
        
        if text:
            return [Document(
                page_content=text,
                metadata={"source": url, "type": "website"}
            )]
    except Exception as e:
        raise Exception(f"Failed to load website: {str(e)}")
    return []


def load_youtube(url: str, api_key: str = None) -> List[Document]:
    video_id = extract_youtube_video_id(url)
    if not video_id:
        raise ValueError("Could not extract YouTube video ID from URL")
    
    text = None
    video_title = ""
    
    if api_key:
        try:
            text, video_title = _fetch_youtube_via_gemini(url, api_key)
        except Exception as e:
            print(f"Gemini 방식 실패: {e}")
    
    if not text:
        try:
            ytt_api = YouTubeTranscriptApi()
            transcript_list = ytt_api.list(video_id)
            
            transcript = None
            for lang_codes in [['ko'], ['en'], ['ko-auto', 'en-auto']]:
                try:
                    transcript = transcript_list.find_transcript(lang_codes)
                    break
                except Exception:
                    continue
            
            if transcript is None:
                available = list(transcript_list)
                if available:
                    transcript = available[0]
            
            if transcript:
                transcript_data = transcript.fetch()
                text = " ".join([snippet.text for snippet in transcript_data])
        except Exception as e:
            print(f"Transcript API 방식 실패: {e}")
    
    if not text or len(text.strip()) < 10:
        raise Exception(
            "YouTube 자막을 가져올 수 없습니다. "
            "이 동영상에 자막이 없거나 YouTube가 클라우드 서버 접근을 차단했을 수 있습니다."
        )
    
    return [Document(
        page_content=text,
        metadata={
            "source": url,
            "video_id": video_id,
            "type": "youtube",
            "title": video_title,
            "filename": f"[YouTube] {video_title[:50]}" if video_title else f"[YouTube] {video_id}"
        }
    )]


def _fetch_youtube_via_gemini(url: str, api_key: str) -> tuple:
    from google import genai
    
    client = genai.Client(api_key=api_key)
    
    response = client.models.generate_content(
        model="gemini-2.0-flash",
        contents=[
            {
                "parts": [
                    {"text": "이 YouTube 동영상의 전체 내용을 한국어로 상세하게 정리해 주세요. "
                             "동영상의 제목, 발표자/출연자, 주요 내용, 핵심 포인트를 모두 포함하여 "
                             "가능한 한 상세하고 빠짐없이 텍스트로 변환해 주세요. "
                             "응답 형식: 첫 줄에 '제목: [동영상 제목]'을 적고, "
                             "그 다음부터 전체 내용을 정리해 주세요."},
                    {"file_data": {"file_uri": url, "mime_type": "video/*"}}
                ]
            }
        ]
    )
    
    content = response.text
    if not content:
        raise Exception("Gemini에서 응답을 받지 못했습니다.")
    
    video_title = ""
    lines = content.split('\n')
    for line in lines:
        if line.strip().startswith('제목:'):
            video_title = line.strip().replace('제목:', '').strip()
            break
    
    return content, video_title


def load_file(uploaded_file) -> List[Document]:
    file_extension = uploaded_file.name.lower().split(".")[-1]
    
    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_extension}") as tmp_file:
        tmp_file.write(uploaded_file.getvalue())
        tmp_path = tmp_file.name
    
    try:
        if file_extension == "pdf":
            documents = load_pdf(tmp_path)
        elif file_extension in ["doc", "docx"]:
            documents = load_docx(tmp_path)
        elif file_extension in ["ppt", "pptx"]:
            documents = load_pptx(tmp_path)
        elif file_extension in ["xls", "xlsx"]:
            documents = load_xlsx(tmp_path)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")
        
        for doc in documents:
            doc.metadata["filename"] = uploaded_file.name
        
        return documents
    finally:
        os.unlink(tmp_path)


def load_url(url: str, api_key: str = None) -> List[Document]:
    if "youtube.com" in url or "youtu.be" in url:
        return load_youtube(url, api_key=api_key)
    else:
        return load_website(url)


def search_google_news(query: str, search_date: str, max_results: int = 20) -> List[Document]:
    """Search Google News for articles about a query for a given month.
    
    Args:
        query: Search query string
        search_date: Date string in 'YYYY-MM' format (month-based search)
        max_results: Maximum number of news articles to fetch
    
    Returns:
        List of Document objects with type='news' metadata
    """
    from datetime import datetime, timedelta
    from urllib.parse import quote_plus
    import calendar
    
    if len(search_date) == 7:
        year, month = search_date.split('-')
        year, month = int(year), int(month)
        date_after = f"{year}-{month:02d}-01"
        last_day = calendar.monthrange(year, month)[1]
        date_before = f"{year}-{month:02d}-{last_day}"
        date_min_fmt = f"{month:02d}/01/{year}"
        date_max_fmt = f"{month:02d}/{last_day}/{year}"
    else:
        date_obj = datetime.strptime(search_date, '%Y-%m-%d')
        date_after = date_obj.strftime('%Y-%m-%d')
        date_before = (date_obj + timedelta(days=1)).strftime('%Y-%m-%d')
        date_min_fmt = date_obj.strftime('%m/%d/%Y')
        date_max_fmt = date_obj.strftime('%m/%d/%Y')
    
    encoded_query = quote_plus(query)
    
    search_urls = [
        f"https://www.google.com/search?q={encoded_query}&tbm=nws&tbs=cdr:1,cd_min:{date_min_fmt},cd_max:{date_max_fmt}",
        f"https://news.google.com/rss/search?q={encoded_query}+after:{date_after}+before:{date_before}&hl=ko&gl=KR&ceid=KR:ko",
    ]
    
    import random
    user_agents = [
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/131.0.0.0 Safari/537.36",
        "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/131.0.0.0 Safari/537.36",
        "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/131.0.0.0 Safari/537.36",
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64; rv:133.0) Gecko/20100101 Firefox/133.0",
    ]
    headers = {
        "User-Agent": random.choice(user_agents),
        "Accept-Language": "ko-KR,ko;q=0.9,en-US;q=0.8,en;q=0.7",
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
    }
    
    documents = []
    visited_urls = set()
    
    try:
        rss_url = search_urls[1]
        print(f"[News Search] RSS URL: {rss_url}")
        response = None
        for attempt in range(3):
            try:
                response = requests.get(rss_url, headers={**headers, "User-Agent": random.choice(user_agents)}, timeout=30)
                if response.status_code == 200:
                    break
                print(f"[News Search] RSS attempt {attempt+1} status: {response.status_code}")
            except Exception as retry_e:
                print(f"[News Search] RSS attempt {attempt+1} error: {retry_e}")
                import time as _time
                _time.sleep(1)
        if response is None:
            raise Exception("All RSS retry attempts failed")
        print(f"[News Search] RSS status: {response.status_code}, content length: {len(response.content)}")
        if response.status_code == 200:
            import xml.etree.ElementTree as ET
            root = ET.fromstring(response.content)
            
            for item in root.iter('item'):
                if len(documents) >= max_results:
                    break
                    
                title = item.find('title')
                link = item.find('link')
                pub_date = item.find('pubDate')
                description = item.find('description')
                
                title_text = title.text if title is not None else ""
                link_text = link.text if link is not None else ""
                pub_date_text = pub_date.text if pub_date is not None else ""
                desc_text = description.text if description is not None else ""
                
                if desc_text:
                    desc_soup = BeautifulSoup(desc_text, "html.parser")
                    desc_text = desc_soup.get_text(strip=True)
                
                if link_text and link_text not in visited_urls:
                    visited_urls.add(link_text)
                    
                    full_content = fetch_news_article(link_text, headers)
                    
                    if full_content and len(full_content.strip()) > 100:
                        content = full_content
                    else:
                        content = f"제목: {title_text}\n\n{desc_text}" if desc_text else f"제목: {title_text}"
                    
                    if content and len(content.strip()) > 50:
                        documents.append(Document(
                            page_content=content,
                            metadata={
                                "source": link_text,
                                "title": title_text,
                                "type": "news",
                                "filename": f"[뉴스] {title_text[:50]}",
                                "search_date": search_date,
                                "search_query": query,
                                "published_date": pub_date_text,
                            }
                        ))
    except Exception as e:
        print(f"[News Search] RSS search error: {e}")
        import traceback
        traceback.print_exc()
    
    print(f"[News Search] After RSS: {len(documents)} documents found")
    
    if len(documents) < 3:
        try:
            google_url = search_urls[0]
            print(f"[News Search] Google URL: {google_url}")
            response = requests.get(google_url, headers=headers, timeout=30)
            if response.status_code == 200:
                soup = BeautifulSoup(response.text, "html.parser")
                
                for article in soup.select('div.SoaBEf, div.dbsr, div.g'):
                    if len(documents) >= max_results:
                        break
                    
                    link_tag = article.find('a', href=True)
                    if not link_tag:
                        continue
                    
                    href = link_tag.get('href', '')
                    if href.startswith('/url?q='):
                        href = href.split('/url?q=')[1].split('&')[0]
                    
                    if not href.startswith('http') or href in visited_urls:
                        continue
                    
                    visited_urls.add(href)
                    
                    title_el = article.find('div', class_='BNeawe') or article.find('h3')
                    title_text = title_el.get_text(strip=True) if title_el else ""
                    
                    snippet_el = article.find('div', class_='BNeawe s3v9rd')
                    snippet_text = snippet_el.get_text(strip=True) if snippet_el else ""
                    
                    full_content = fetch_news_article(href, headers)
                    
                    if full_content and len(full_content.strip()) > 100:
                        content = full_content
                    else:
                        content = f"제목: {title_text}\n\n{snippet_text}" if snippet_text else f"제목: {title_text}"
                    
                    if content and len(content.strip()) > 50:
                        documents.append(Document(
                            page_content=content,
                            metadata={
                                "source": href,
                                "title": title_text,
                                "type": "news",
                                "filename": f"[뉴스] {title_text[:50]}",
                                "search_date": search_date,
                                "search_query": query,
                            }
                        ))
        except Exception as e:
            print(f"[News Search] Google search error: {e}")
            import traceback
            traceback.print_exc()
    
    print(f"[News Search] Total documents found: {len(documents)}")
    return documents


def fetch_news_article(url: str, headers: dict) -> Optional[str]:
    """Fetch and extract the main text content from a news article URL."""
    try:
        if any(blocked in url for blocked in ['google.com/url', 'accounts.google', 'consent.google']):
            return None
            
        response = requests.get(url, headers=headers, timeout=10, allow_redirects=True)
        response.raise_for_status()
        
        content_type = response.headers.get('Content-Type', '')
        if 'text/html' not in content_type:
            return None
        
        soup = BeautifulSoup(response.text, "html.parser")
        
        for tag in soup(["script", "style", "nav", "footer", "header", "aside", "iframe", "noscript", "form"]):
            tag.decompose()
        
        article = soup.find('article') or soup.find('div', class_=re.compile(r'article|content|body|post', re.I))
        
        if article:
            text = article.get_text(separator="\n", strip=True)
        else:
            text = soup.get_text(separator="\n", strip=True)
        
        text = re.sub(r'\n\s*\n', '\n\n', text)
        
        lines = text.split('\n')
        filtered_lines = [line for line in lines if len(line.strip()) > 10]
        text = '\n'.join(filtered_lines)
        
        if len(text) > 5000:
            text = text[:5000]
        
        return text
    except Exception:
        return None


def split_documents(documents: List[Document], chunk_size: int = 2000, chunk_overlap: int = 400) -> List[Document]:
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
        separators=["\n\n", "\n", "。", ". ", ", ", " ", ""],
        is_separator_regex=False
    )
    return text_splitter.split_documents(documents)
